# -*- coding: utf-8 -*-
"""
将资料库内的 PPTX / DOCX 转换为自包含、移动端友好的在线网页（HTML）。
- 抽取每页文字与图片；图片保存为同目录 _files/ 文件夹并相对引用，PIL 压缩到最大宽1400px。
- 生成的 .html 与原件同目录（兄弟文件），原件保留作「下载原文件」按钮。
- 用法：python gen_doc_viewers.py
"""
import os, mimetypes
from io import BytesIO
from pptx import Presentation
from docx import Document
import xml.etree.ElementTree as ET
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))

TARGETS = [
    ("files/公司介绍/2604新华在售产品培训.pptx", "ppt"),
    ("files/公司介绍/新华保险投资实力介绍（新华资管）.pptx", "ppt"),
    ("files/产品对比/招行在架产品对比分析.pptx", "ppt"),
    ("files/产品规则/荣尊世家减保规则.pptx", "ppt"),
    ("files/产品介绍/新华宏御世家_知识问答卡片 (6).docx", "doc"),
]

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

def html_escape(s):
    return (s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;'))

def rel_path_to_index(html_path):
    rel = os.path.relpath(os.path.join(ROOT, "index.html"), os.path.dirname(html_path))
    return rel.replace('\\', '/')

class ImgWriter:
    def __init__(self, html_path):
        self.dir = os.path.splitext(html_path)[0] + '_files'
        self.base_rel = os.path.basename(self.dir)
        self.n = 0
        os.makedirs(self.dir, exist_ok=True)
    def add(self, blob, ext):
        self.n += 1
        try:
            im = Image.open(BytesIO(blob))
            maxw = 1400
            if im.width > maxw:
                ratio = maxw / im.width
                im = im.resize((maxw, int(im.height * ratio)), Image.LANCZOS)
            buf = BytesIO()
            if im.mode in ('RGBA', 'LA', 'P'):
                ext_out = 'png'; im.save(buf, 'PNG')
            else:
                im = im.convert('RGB'); im.save(buf, 'JPEG', quality=82)
            data = buf.getvalue(); ext_out = ext_out
        except Exception:
            data = blob; ext_out = (ext or 'png').lower().lstrip('.')
        fname = f'img{self.n}.{ext_out}'
        with open(os.path.join(self.dir, fname), 'wb') as f:
            f.write(data)
        return f'{self.base_rel}/{fname}'

# ---------------- PPTX ----------------
def pptx_slide_to_blocks(slide, iw):
    texts, images = [], []
    for shape in slide.shapes:
        if shape.shape_type == 13:
            try:
                images.append(f'<div class="slide-img"><img src="{iw.add(shape.image.blob, shape.image.ext)}" alt=""></div>')
            except Exception:
                pass
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            txt = ''.join(r.text for r in para.runs).strip()
            if not txt:
                continue
            size = None
            for r in para.runs:
                if r.font.size:
                    size = r.font.size.pt; break
            texts.append((txt, size))
    if not texts:
        return '', [], images
    texts_sorted = sorted(texts, key=lambda t: (t[1] if t[1] else 0), reverse=True)
    title = texts_sorted[0][0]
    others = [t[0] for t in texts if t[0] != title]
    return title, others, images

def convert_pptx(src, out_html, original_rel):
    iw = ImgWriter(out_html)
    prs = Presentation(src)
    slides_html = []
    for i, slide in enumerate(prs.slides, 1):
        title, bullets, images = pptx_slide_to_blocks(slide, iw)
        body = ''
        if bullets:
            lis = ''.join(f'<li>{html_escape(b)}</li>' for b in bullets)
            body += f'<ul class="bullets">{lis}</ul>'
        imgs = ''.join(images)
        slides_html.append(f'''
        <section class="slide" data-idx="{i}">
          <div class="slide-inner">
            {f'<h2 class="slide-title">{html_escape(title)}</h2>' if title else ''}
            {body}
            {imgs}
          </div>
        </section>''')
    total = len(slides_html)
    html = TEMPLATE_DECK.format(
        title=html_escape(os.path.basename(src)), total=total,
        slides='\n'.join(slides_html), original=original_rel,
        index=rel_path_to_index(out_html), dlname=html_escape(os.path.basename(src)),
    )
    with open(out_html, 'w', encoding='utf-8') as f:
        f.write(html)
    return total

# ---------------- DOCX ----------------
def convert_docx(src, out_html, original_rel):
    iw = ImgWriter(out_html)
    doc = Document(src)
    parts = doc.part
    blocks = []
    for el in doc.element.body.iterchildren():
        if el.tag.endswith('}p'):
            from docx.text.paragraph import Paragraph
            para = Paragraph(el, doc)
            text = para.text.strip()
            imgs = []
            for run in para.runs:
                blip = run._element.find('.//a:blip', NS)
                if blip is not None:
                    rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if rId and rId in parts.related_parts:
                        rel = parts.related_parts[rId]
                        imgs.append(f'<div class="doc-img"><img src="{iw.add(rel.blob, os.path.splitext(str(rel.partname))[1].lstrip(".") or "png")}" alt=""></div>')
            style = (para.style.name or '') if para.style else ''
            if not text and not imgs:
                continue
            if imgs and not text:
                blocks.append(''.join(imgs)); continue
            cls, lvl = 'doc-p', 0
            if 'Heading' in style or '标题' in style:
                try: lvl = int(''.join(filter(str.isdigit, style)) or '1')
                except Exception: lvl = 1
                cls = f'doc-h{lvl}' if lvl <= 3 else 'doc-h3'
            esc = html_escape(text)
            blocks.append((f'<p class="{cls}">{esc}</p>' + ''.join(imgs)) if imgs else f'<p class="{cls}">{esc}</p>')
        elif el.tag.endswith('}tbl'):
            from docx.table import Table
            table = Table(el, doc)
            rows = ''.join('<tr>' + ''.join(f'<td>{html_escape(c.text.strip())}</td>' for c in r.cells) + '</tr>' for r in table.rows)
            blocks.append(f'<div class="doc-table-wrap"><table class="doc-table">{rows}</table></div>')
    body = '\n'.join(blocks)
    html = TEMPLATE_DOC.format(
        title=html_escape(os.path.basename(src)), body=body,
        original=original_rel, index=rel_path_to_index(out_html),
        dlname=html_escape(os.path.basename(src)),
    )
    with open(out_html, 'w', encoding='utf-8') as f:
        f.write(html)
    return len(blocks)

# ---------------- 模板 ----------------
TEMPLATE_DECK = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
<meta name="apple-mobile-web-app-capable" content="yes">
<title>{title} · 在线浏览</title>
<style>
  *{{margin:0;padding:0;box-sizing:border-box;-webkit-tap-highlight-color:transparent;}}
  body{{font-family:-apple-system,BlinkMacSystemFont,'PingFang SC','Microsoft YaHei',sans-serif;background:#0F172A;color:#1E293B;overflow:hidden;}}
  .topbar{{position:fixed;top:0;left:0;right:0;z-index:30;display:flex;align-items:center;gap:10px;padding:10px 14px;background:linear-gradient(135deg,#1D4ED8,#2563EB);color:#fff;}}
  .topbar a.back{{color:#fff;text-decoration:none;font-size:14px;white-space:nowrap;}}
  .topbar .t{{flex:1;font-size:14px;font-weight:600;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;}}
  .topbar .cnt{{font-size:13px;opacity:.9;white-space:nowrap;}}
  .stage{{position:absolute;top:52px;left:0;right:0;bottom:96px;overflow:hidden;}}
  .slides{{display:flex;height:100%;transition:transform .35s ease;}}
  .slide{{min-width:100%;height:100%;display:flex;align-items:center;justify-content:center;padding:14px;}}
  .slide-inner{{width:100%;max-width:760px;background:#fff;border-radius:16px;padding:26px 22px;max-height:100%;overflow-y:auto;-webkit-overflow-scrolling:touch;box-shadow:0 10px 40px rgba(0,0,0,.25);}}
  .slide-title{{font-size:21px;font-weight:800;color:#1D4ED8;margin-bottom:14px;line-height:1.4;}}
  .bullets{{list-style:none;}}
  .bullets li{{position:relative;padding:7px 0 7px 20px;font-size:16px;line-height:1.6;border-bottom:1px solid #F1F5F9;}}
  .bullets li:before{{content:'';position:absolute;left:4px;top:16px;width:7px;height:7px;border-radius:50%;background:#2563EB;}}
  .slide-img{{margin:12px 0;}}
  .slide-img img{{width:100%;border-radius:10px;display:block;}}
  .dots{{position:fixed;left:0;right:0;bottom:54px;z-index:30;display:flex;justify-content:center;gap:6px;flex-wrap:wrap;padding:0 20px;}}
  .dots span{{width:7px;height:7px;border-radius:50%;background:rgba(255,255,255,.4);transition:all .2s;}}
  .dots span.on{{background:#fff;width:18px;border-radius:4px;}}
  .bottombar{{position:fixed;left:0;right:0;bottom:0;z-index:30;display:flex;gap:10px;padding:10px 16px calc(10px + env(safe-area-inset-bottom));background:rgba(15,23,42,.96);}}
  .bottombar button{{flex:1;padding:11px;border:none;border-radius:10px;font-size:14px;font-weight:600;cursor:pointer;}}
  .bottombar .prev{{background:#334155;color:#fff;}}
  .bottombar .next{{background:linear-gradient(135deg,#2563EB,#1D4ED8);color:#fff;}}
  .bottombar .dl{{background:#16A34A;color:#fff;text-decoration:none;display:flex;align-items:center;justify-content:center;}}
  .bottombar button:active{{transform:scale(.97);}}
</style>
</head>
<body>
  <div class="topbar">
    <a class="back" href="{index}">← 返回</a>
    <div class="t">{title}</div>
    <div class="cnt"><span id="cur">1</span>/{total}</div>
  </div>
  <div class="stage"><div class="slides" id="slides">
{slides}
  </div></div>
  <div class="dots" id="dots"></div>
  <div class="bottombar">
    <button class="prev" onclick="go(-1)">‹ 上一页</button>
    <button class="next" onclick="go(1)">下一页 ›</button>
    <a class="dl" href="{original}" download="{dlname}">⬇ 下载原文件</a>
  </div>
<script>
  const slidesEl=document.getElementById('slides');const total={total};let cur=0;
  const dotsWrap=document.getElementById('dots');
  for(let i=0;i<total;i++){{const s=document.createElement('span');if(i===0)s.className='on';dotsWrap.appendChild(s);}}
  function render(){{slidesEl.style.transform='translateX('+(-cur*100)+'%)';document.getElementById('cur').textContent=cur+1;dotsWrap.querySelectorAll('span').forEach((d,i)=>d.classList.toggle('on',i===cur));}}
  function go(d){{cur=Math.max(0,Math.min(total-1,cur+d));render();}}
  let sx=0,moving=false;
  slidesEl.addEventListener('touchstart',e=>{{sx=e.touches[0].clientX;moving=true;}},{{passive:true}});
  slidesEl.addEventListener('touchend',e=>{{if(!moving)return;moving=false;const dx=e.changedTouches[0].clientX-sx;if(Math.abs(dx)>50)go(dx<0?1:-1);}},{{passive:true}});
  render();
</script>
</body>
</html>"""

TEMPLATE_DOC = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
<meta name="apple-mobile-web-app-capable" content="yes">
<title>{title} · 在线浏览</title>
<style>
  *{{margin:0;padding:0;box-sizing:border-box;-webkit-tap-highlight-color:transparent;}}
  body{{font-family:-apple-system,BlinkMacSystemFont,'PingFang SC','Microsoft YaHei',sans-serif;background:#F1F5F9;color:#1E293B;line-height:1.7;}}
  .topbar{{position:sticky;top:0;z-index:30;display:flex;align-items:center;gap:10px;padding:10px 14px;background:linear-gradient(135deg,#1D4ED8,#2563EB);color:#fff;}}
  .topbar a.back{{color:#fff;text-decoration:none;font-size:14px;white-space:nowrap;}}
  .topbar .t{{flex:1;font-size:14px;font-weight:600;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;}}
  .topbar .dl{{background:#16A34A;color:#fff;text-decoration:none;font-size:13px;padding:7px 12px;border-radius:8px;white-space:nowrap;}}
  .doc{{max-width:720px;margin:0 auto;padding:20px 18px 60px;}}
  .doc-h1{{font-size:23px;font-weight:800;color:#1D4ED8;margin:18px 0 12px;}}
  .doc-h2{{font-size:19px;font-weight:700;color:#1E40AF;margin:16px 0 8px;}}
  .doc-h3{{font-size:16px;font-weight:700;color:#334155;margin:12px 0 6px;}}
  .doc-p{{font-size:15px;margin:8px 0;color:#334155;}}
  .doc-img{{margin:14px 0;text-align:center;}}
  .doc-img img{{max-width:100%;border-radius:10px;box-shadow:0 4px 16px rgba(0,0,0,.1);}}
  .doc-table-wrap{{overflow-x:auto;margin:12px 0;-webkit-overflow-scrolling:touch;}}
  .doc-table{{border-collapse:collapse;width:100%;font-size:13px;background:#fff;}}
  .doc-table td{{border:1px solid #E2E8F0;padding:7px 9px;}}
  .doc-table tr:nth-child(even) td{{background:#F8FAFC;}}
</style>
</head>
<body>
  <div class="topbar">
    <a class="back" href="{index}">← 返回</a>
    <div class="t">{title}</div>
    <a class="dl" href="{original}" download="{dlname}">⬇ 下载原文件</a>
  </div>
  <div class="doc">
{body}
  </div>
</body>
</html>"""

def main():
    # 清理旧的生成文件，避免残留
    for src_rel, _ in TARGETS:
        stem, _ = os.path.splitext(os.path.join(ROOT, src_rel))
        for p in (stem + '.html', stem + '_files'):
            if os.path.isfile(p): os.remove(p)
            elif os.path.isdir(p):
                import shutil; shutil.rmtree(p)
    for src_rel, kind in TARGETS:
        src = os.path.join(ROOT, src_rel)
        if not os.path.exists(src):
            print(f'[跳过] 不存在: {src_rel}'); continue
        stem, _ = os.path.splitext(src)
        out_html = stem + '.html'
        original_rel = os.path.basename(src)
        if kind == 'ppt':
            n = convert_pptx(src, out_html, original_rel)
            print(f'[PPT] {src_rel} -> {os.path.basename(out_html)} ({n} 页)')
        else:
            n = convert_docx(src, out_html, original_rel)
            print(f'[DOC] {src_rel} -> {os.path.basename(out_html)} ({n} 块)')

if __name__ == '__main__':
    main()
